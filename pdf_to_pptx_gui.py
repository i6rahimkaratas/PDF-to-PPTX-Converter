import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os
import sys
from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io
import tempfile
import threading
import shutil

class PDFToPPTXConverter:
    def __init__(self):
        """PDF'den PPTX'e dönüştürücü sınıfı"""
        self.temp_dir = tempfile.mkdtemp()
        
    def extract_text_and_images_from_pdf(self, pdf_path, progress_callback=None):
        """PDF'den metin ve görselleri çıkarır"""
        doc = fitz.open(pdf_path)
        pages_data = []
        total_pages = len(doc)
        
        for page_num in range(total_pages):
            page = doc[page_num]
            page_data = {
                'page_number': page_num + 1,
                'text': '',
                'images': []
            }
            
            # Metni çıkar
            text = page.get_text()
            page_data['text'] = text.strip()
            
            # Görselleri çıkar
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Geçici dosya olarak kaydet
                    image_filename = f"page_{page_num + 1}_img_{img_index + 1}.{image_ext}"
                    image_path = os.path.join(self.temp_dir, image_filename)
                    
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    page_data['images'].append(image_path)
                except Exception as e:
                    print(f"Görsel çıkarma hatası: {e}")
            
            pages_data.append(page_data)
            
            # İlerleme güncellemesi
            if progress_callback:
                progress = (page_num + 1) / total_pages * 50  # İlk %50 PDF okuma
                progress_callback(progress, f"Sayfa {page_num + 1}/{total_pages} işleniyor...")
        
        doc.close()
        return pages_data
    
    def create_slide_with_content(self, prs, page_data):
        """Sayfa verisi ile yeni bir slayt oluşturur"""
        # Boş slayt düzeni
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Slayt boyutları
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Başlık ekle
        if page_data['text']:
            # İlk satırı başlık olarak kullan
            lines = page_data['text'].split('\n')
            title_text = lines[0] if lines and lines[0].strip() else f"Sayfa {page_data['page_number']}"
            
            # Başlık text box'ı
            title_left = Inches(0.5)
            title_top = Inches(0.5)
            title_width = slide_width - Inches(1)
            title_height = Inches(1)
            
            title_textbox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_textbox.text_frame
            title_frame.text = title_text[:100]  # Başlığı sınırla
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Geri kalan metni içerik olarak ekle
            if len(lines) > 1:
                content_text = '\n'.join(lines[1:])
                
                content_left = Inches(0.5)
                content_top = Inches(1.8)
                content_width = slide_width - Inches(1)
                
                # Görsel varsa metin alanını küçült
                if page_data['images']:
                    content_height = Inches(3)
                else:
                    content_height = slide_height - Inches(2.5)
                
                content_textbox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                content_frame = content_textbox.text_frame
                content_frame.text = content_text[:2000]  # Metni sınırla
                content_frame.word_wrap = True
                
                for paragraph in content_frame.paragraphs:
                    paragraph.font.size = Pt(12)
        
        # Görselleri ekle
        if page_data['images']:
            img_count = len(page_data['images'])
            img_width = Inches(2.5)
            img_height = Inches(2)
            
            # Görselleri yan yana yerleştir
            start_left = (slide_width - (img_width * min(img_count, 3))) / 2
            
            for i, img_path in enumerate(page_data['images'][:3]):  # Maksimum 3 görsel
                try:
                    img_left = start_left + (img_width * i)
                    img_top = slide_height - img_height - Inches(0.5)
                    
                    slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)
                except Exception as e:
                    print(f"Görsel eklenirken hata: {e}")
    
    def convert_pdf_to_pptx(self, pdf_path, output_path, progress_callback=None):
        """PDF'yi PPTX'e dönüştürür"""
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF dosyası bulunamadı: {pdf_path}")
        
        # PDF'den veri çıkar
        pages_data = self.extract_text_and_images_from_pdf(pdf_path, progress_callback)
        
        if progress_callback:
            progress_callback(60, "PowerPoint sunumu oluşturuluyor...")
        
        # PowerPoint sunumu oluştur
        prs = Presentation()
        
        # Her sayfa için slayt oluştur
        total_pages = len(pages_data)
        for i, page_data in enumerate(pages_data):
            self.create_slide_with_content(prs, page_data)
            
            if progress_callback:
                progress = 60 + (i + 1) / total_pages * 35  # %60-95 arası slayt oluşturma
                progress_callback(progress, f"Slayt {i + 1}/{total_pages} oluşturuluyor...")
        
        # İlk boş slaytı sil (eğer varsa)
        if len(prs.slides) > len(pages_data):
            try:
                slide_id = prs.slides._sldIdLst[0]
                prs.part.drop_rel(slide_id.rId)
                prs.slides._sldIdLst.remove(slide_id)
            except:
                pass
        
        if progress_callback:
            progress_callback(95, "Dosya kaydediliyor...")
        
        # PPTX dosyasını kaydet
        prs.save(output_path)
        
        # Geçici dosyaları temizle
        self.cleanup_temp_files()
        
        if progress_callback:
            progress_callback(100, "Tamamlandı!")
        
        return output_path
    
    def cleanup_temp_files(self):
        """Geçici dosyaları temizler"""
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)

class PDFToPPTXGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF'den PPTX'e Dönüştürücü")
        self.root.geometry("600x500")
        self.root.resizable(True, True)
        
        # Değişkenler
        self.pdf_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.converter = None
        
        self.setup_ui()
        
        # Pencereyi merkeze al
        self.center_window()
    
    def center_window(self):
        """Pencereyi ekranın merkezine alır"""
        self.root.update_idletasks()
        width = self.root.winfo_width()
        height = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (width // 2)
        y = (self.root.winfo_screenheight() // 2) - (height // 2)
        self.root.geometry(f'{width}x{height}+{x}+{y}')
    
    def setup_ui(self):
        """Kullanıcı arayüzünü oluşturur"""
        # Ana çerçeve
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # Başlık
        title_label = ttk.Label(main_frame, text="PDF'den PowerPoint'e Dönüştürücü", 
                               font=('Arial', 16, 'bold'))
        title_label.grid(row=0, column=0, columnspan=3, pady=(0, 20))
        
        # PDF dosya seçimi
        ttk.Label(main_frame, text="PDF Dosyası:").grid(row=1, column=0, sticky=tk.W, pady=5)
        
        pdf_frame = ttk.Frame(main_frame)
        pdf_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        
        ttk.Entry(pdf_frame, textvariable=self.pdf_path, width=50).pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(pdf_frame, text="Gözat", command=self.browse_pdf).pack(side=tk.RIGHT, padx=(10, 0))
        
        # Çıktı dosya seçimi
        ttk.Label(main_frame, text="Kaydet:").grid(row=3, column=0, sticky=tk.W, pady=(10, 5))
        
        output_frame = ttk.Frame(main_frame)
        output_frame.grid(row=4, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 20))
        
        ttk.Entry(output_frame, textvariable=self.output_path, width=50).pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(output_frame, text="Konum Seç", command=self.browse_output).pack(side=tk.RIGHT, padx=(10, 0))
        
        # İlerleme çubuğu
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(main_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=10)
        
        # Durum etiketi
        self.status_var = tk.StringVar(value="Dosyaları seçin ve dönüştürmeye başlayın")
        self.status_label = ttk.Label(main_frame, textvariable=self.status_var, foreground="blue")
        self.status_label.grid(row=6, column=0, columnspan=3, pady=5)
        
        # Dönüştür butonu
        self.convert_button = ttk.Button(main_frame, text="Dönüştür", command=self.start_conversion)
        self.convert_button.grid(row=7, column=0, columnspan=3, pady=20)
        
        # Bilgi kutusu
        info_frame = ttk.LabelFrame(main_frame, text="Bilgi", padding="10")
        info_frame.grid(row=8, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=10)
        
        info_text = """• PDF dosyanızı seçin
• PPTX dosyasının kaydedileceği konumu belirleyin
• Dönüştür butonuna tıklayın
• İşlem tamamlandığında dosyanız belirtilen konuma kaydedilecektir"""
        
        ttk.Label(info_frame, text=info_text, justify=tk.LEFT).pack(anchor=tk.W)
        
        # Grid ağırlıkları
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        pdf_frame.columnconfigure(0, weight=1)
        output_frame.columnconfigure(0, weight=1)
    
    def browse_pdf(self):
        """PDF dosyası seçme dialog'u"""
        filename = filedialog.askopenfilename(
            title="PDF Dosyası Seçin",
            filetypes=[("PDF dosyaları", "*.pdf"), ("Tüm dosyalar", "*.*")]
        )
        if filename:
            self.pdf_path.set(filename)
            # Otomatik çıktı yolu öner
            suggested_output = filename.replace('.pdf', '_converted.pptx')
            self.output_path.set(suggested_output)
    
    def browse_output(self):
        """Çıktı dosyası kaydetme dialog'u"""
        filename = filedialog.asksaveasfilename(
            title="PPTX Dosyasını Kaydet",
            defaultextension=".pptx",
            filetypes=[("PowerPoint dosyaları", "*.pptx"), ("Tüm dosyalar", "*.*")]
        )
        if filename:
            self.output_path.set(filename)
    
    def update_progress(self, value, message):
        """İlerleme çubuğunu günceller"""
        self.progress_var.set(value)
        self.status_var.set(message)
        self.root.update_idletasks()
    
    def start_conversion(self):
        """Dönüştürme işlemini başlatır"""
        if not self.pdf_path.get():
            messagebox.showerror("Hata", "Lütfen bir PDF dosyası seçin!")
            return
        
        if not self.output_path.get():
            messagebox.showerror("Hata", "Lütfen çıktı dosyası konumunu belirleyin!")
            return
        
        if not os.path.exists(self.pdf_path.get()):
            messagebox.showerror("Hata", "Seçilen PDF dosyası bulunamadı!")
            return
        
        # Butonu deaktif et
        self.convert_button.config(state="disabled")
        self.progress_var.set(0)
        
        # Dönüştürme işlemini ayrı thread'de çalıştır
        thread = threading.Thread(target=self.convert_file)
        thread.daemon = True
        thread.start()
    
    def convert_file(self):
        """Dosya dönüştürme işlemi"""
        try:
            self.converter = PDFToPPTXConverter()
            
            result_path = self.converter.convert_pdf_to_pptx(
                self.pdf_path.get(),
                self.output_path.get(),
                self.update_progress
            )
            
            # Başarılı mesajı
            self.root.after(0, lambda: messagebox.showinfo(
                "Başarılı!", 
                f"Dönüştürme tamamlandı!\n\nDosya kaydedildi:\n{result_path}"
            ))
            
        except Exception as e:
            # Hata mesajı
            self.root.after(0, lambda: messagebox.showerror(
                "Hata",
                f"Dönüştürme sırasında hata oluştu:\n{str(e)}"
            ))
            self.root.after(0, lambda: self.status_var.set("Hata oluştu!"))
        
        finally:
            # Butonu tekrar aktif et
            self.root.after(0, lambda: self.convert_button.config(state="normal"))

def main():
    """Ana fonksiyon"""
    try:
        # Gerekli kütüphanelerin kontrolü
        import fitz
        from pptx import Presentation
        from PIL import Image
    except ImportError as e:
        print(f"Gerekli kütüphane eksik: {e}")
        print("Lütfen şu komutları çalıştırın:")
        print("pip install PyMuPDF python-pptx Pillow")
        return
    
    # GUI başlat
    root = tk.Tk()
    app = PDFToPPTXGUI(root)
    root.mainloop()

if __name__ == "__main__":
    main()